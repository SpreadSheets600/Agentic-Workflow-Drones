"""Generate the 7-slide challenge deck: WeevilDrone_Agentic_Demo.pptx
Run:  .venv/bin/python scripts/make_slides.py
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0F, 0x2A, 0x43)
TEAL = RGBColor(0x14, 0x8F, 0x87)
ORANGE = RGBColor(0xE8, 0x8B, 0x1A)
RED = RGBColor(0xC0, 0x45, 0x3E)
GRAY = RGBColor(0x44, 0x4C, 0x54)
LIGHT = RGBColor(0xEF, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MONO = "Consolas"
HEAD = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(12.4), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HEAD
    if subtitle:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xBF, 0xD4, 0xD4); r2.font.name = HEAD
    return s


def bullets(s, items, x=0.6, y=1.25, w=12.1, h=5.9, size=17):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        text, lvl, opts = (item + ({},))[:3] if isinstance(item, tuple) else (item, 0, {})
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(opts.get("space", 8))
        r = p.add_run(); r.text = text
        r.font.size = Pt(opts.get("size", size if lvl == 0 else size - 2))
        r.font.color.rgb = opts.get("color", GRAY)
        r.font.bold = opts.get("bold", False)
        r.font.name = MONO if opts.get("mono") else HEAD
    return tb


def box(s, x, y, w, h, text, fill, font_color=WHITE, size=12, bold=True,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = WHITE; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = font_color; r.font.name = HEAD
    return sh


def arrow(s, x1, y1, x2, y2, color=NAVY, width=2.5, dashed=False):
    conn = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash",
                                 {"val": "dash"}))
    ln.append(ln.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd",
                             {"type": "triangle", "w": "med", "len": "med"}))
    return conn


# ------------------------------------------------ slide 1: problem/mission ---
s = slide("Solar Panel Inspection with an Agentic Drone",
          "WeevilDrone Technology Agentic AI Workflow Challenge")
box(s, 0.6, 1.3, 12.1, 0.95,
    'MISSION:  "Inspect Area A, identify a possible anomaly,\ncollect evidence, and generate a short inspection report."',
    NAVY, WHITE, 16)
bullets(s, [
    ("Why agentic instead of a fixed script?", 0, {"bold": True, "color": NAVY, "size": 19}),
    ("The situation changes mid-mission: cameras fail, evidence is weak, battery drains.", 0),
    ("A fixed workflow always executes the same sequence; an agent decides what to do next from observations.", 0),
    ("What the system demonstrates", 0, {"bold": True, "color": NAVY, "size": 19, "space": 4}),
    ("Goal > Plan > Observe > Decide > Act > Verify > Re-plan > Complete", 0, {"mono": True, "size": 15}),
    ("6 tools · decisions driven by tool results · visible state changes · bounded retry · safe abort", 0),
    ("Two mission intake paths", 0, {"bold": True, "color": NAVY, "size": 19, "space": 4}),
    ("Structured (--scenario) or free text (--mission \"Inspect the north solar field; return when battery hits 45%\")", 0),
    ("The mock-NLP parser proposes target/capabilities/constraints pre-flight still gates every one of them", 0),
    ("Engineering constraints I set", 0, {"bold": True, "color": NAVY, "size": 19, "space": 4}),
    ("Mocked drone + vision (same contracts a real stack would expose) · plain Python · every decision explainable to a line of code", 0),
], y=2.45, h=4.8)

# ------------------------------------------------- slide 2: architecture ---
s = slide("Architecture", "Agent decides WHAT it wants deterministic guardrails decide what it MAY")
box(s, 0.4, 2.6, 1.6, 0.9, "MISSION\nINPUT", GRAY, WHITE, 12)
box(s, 2.4, 1.35, 4.6, 3.9, "", NAVY)
tb = s.shapes.add_textbox(Inches(2.5), Inches(1.42), Inches(4.4), Inches(0.4))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "AGENT CONTROLLER"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
box(s, 2.65, 1.95, 1.9, 0.65, "PLANNER\n(proposal only)", TEAL, WHITE, 11)
box(s, 4.8, 1.95, 1.9, 0.65, "DECISION LOOP\n_decide(): from STATE", TEAL, WHITE, 11)
box(s, 2.65, 2.85, 1.9, 0.65, "POLICIES\nbattery>=30 · conf>=0.70", TEAL, WHITE, 11)
box(s, 4.8, 2.85, 1.9, 0.65, "STATE / MEMORY\nevidence · failures · history", TEAL, WHITE, 11)
box(s, 3.7, 3.9, 3.0, 0.65, "KNOWLEDGE (RAG)\nadvisory, low-confidence only", RGBColor(0x0B, 0x66, 0x60), WHITE, 11)
box(s, 7.6, 2.6, 2.0, 0.9, "GUARDRAILS\ndeterministic checks", ORANGE, WHITE, 12)
box(s, 10.2, 2.6, 2.7, 0.9, "SIMULATED DRONE\nMockDrone · MockVision\nToolObservation out", RGBColor(0x5B, 0x6B, 0x79), WHITE, 11)
box(s, 10.2, 4.6, 2.7, 0.8, "FINAL REPORT\nfrom mission memory", RGBColor(0x5B, 0x6B, 0x79), WHITE, 12)
arrow(s, 2.0, 3.05, 2.4, 3.05)
arrow(s, 4.55, 2.27, 4.8, 2.27)
arrow(s, 5.75, 2.6, 3.6, 2.85)
arrow(s, 6.7, 2.27, 7.6, 2.85)
arrow(s, 9.6, 3.05, 10.2, 3.05)
arrow(s, 11.55, 3.5, 11.55, 4.6, dashed=True)
arrow(s, 10.2, 3.35, 6.7, 2.5, dashed=True, color=TEAL)
bullets(s, [
    ("Decision loop:  AgentDecision (why) > AgentAction (what) > validator > tool > ToolObservation > state update > decide again", 0, {"size": 13, "space": 4}),
    ("Guardrail = pure function (action, state) > (allowed, reason).  The agent cannot bypass or argue with it.", 0, {"size": 13, "space": 4}),
    ("Two memory kinds: mission memory (this flight) drives decisions; RAG knowledge (static docs) is advisory and never executes anything.", 0, {"size": 13, "space": 4}),
    ("Every tool call runs under a watchdog (10s): a hung tool becomes a normal failure observation guardrails still apply.", 0, {"size": 13}),
], y=5.75, h=1.6)

# ---------------------------------------------- slide 3: decision workflow ---
s = slide("Decision Workflow state drives actions, not a plan index")
steps = ["MISSION", "PLAN", "OBSERVE", "DECIDE", "ACT", "VERIFY", "RE-PLAN", "COMPLETE"]
x = 0.45
for i, st in enumerate(steps):
    fill = ORANGE if st in ("DECIDE", "RE-PLAN") else TEAL
    box(s, x, 1.4, 1.62, 0.6, st, fill, WHITE, 12)
    if i < len(steps) - 1:
        arrow(s, x + 1.62, 1.7, x + 1.72, 1.7, width=2)
    x += 1.63
bullets(s, [
    ("Key invariant: _decide() never reads the plan list it reads current state. That is what makes it agentic, not scripted.", 0, {"bold": True, "color": NAVY}),
], y=2.3, h=0.5, size=16)
box(s, 0.6, 2.95, 5.9, 2.7, "", LIGHT)
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(5.5), Inches(0.5))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "Decision point 1 tool failure"
r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = RED
bullets(s, [
    ("Observation: capture_image > success=false 'Camera Timeout'", 0, {"mono": True, "size": 12}),
    ("Interpretation: transient fault > recoverable", 0, {"size": 14}),
    ("Fallback: retry, bounded by a budget of 3 (guardrail-enforced)", 0, {"size": 14}),
    ("4th failure > guardrail rejects capture > safe abort", 0, {"size": 14}),
], x=0.85, y=3.55, w=5.4, h=2.0, size=14)
box(s, 6.9, 2.95, 5.9, 2.7, "", LIGHT)
tb = s.shapes.add_textbox(Inches(7.1), Inches(3.05), Inches(5.5), Inches(0.5))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "Decision point 2 weak evidence"
r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = RED
bullets(s, [
    ("Observation: detect_anomaly > confidence=0.46 (< 0.70)", 0, {"mono": True, "size": 12}),
    ("Interpretation: insufficient evidence is not mission failure", 0, {"size": 14}),
    ("Fallback: consult knowledge base > rewrite plan > re-inspect", 0, {"size": 14}),
    ("Second detection: 0.93 >= 0.70 > verify > return > report", 0, {"size": 14}),
], x=7.15, y=3.55, w=5.4, h=2.0, size=14)
bullets(s, [
    ("Same observation, different decisions chosen by state, budget, and policy. The report prints Initial Plan vs Actual Decision Path as proof of re-planning.", 0, {"bold": True, "color": NAVY, "size": 14}),
], y=5.85, h=0.9, size=14)

# ---------------------------------------------------- slide 4: demo trace ---
s = slide("Demo Trace one mission, end to end",
          "python -m app.main --scenario mission   (all mocked, runs on a laptop)")
trace = [
    ("PRE_FLIGHT", "5/5 capability checks PASS > mission ready", NAVY),
    ("check_battery", "battery=100 > continue", NAVY),
    ("move_to_target", "location BASE > target, battery=90", NAVY),
    ("capture_image", "success=false 'Camera Timeout' > RETRY (1/3)", RED),
    ("capture_image", "success=true evidence=IMG-002", NAVY),
    ("detect_anomaly", "confidence=0.46 < 0.70 > consult RAG > RE-PLAN + RE-INSPECT", ORANGE),
    ("detect_anomaly", "confidence=0.93 >= 0.70 > VERIFY FINDING", NAVY),
    ("return_to_base", "location BASE, battery=80", NAVY),
    ("generate_report", "COMPLETE report generated from mission memory", TEAL),
]
y = 1.3
for tool, result, c in trace:
    box(s, 0.6, y, 2.5, 0.44, tool, c, WHITE, 11)
    tb = s.shapes.add_textbox(Inches(3.3), y, Inches(9.4), Inches(0.44))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = result
    r.font.size = Pt(13); r.font.name = MONO; r.font.color.rgb = c
    y += 0.56
bullets(s, [
    ("Requirements met:  >=3 tool calls OK   decision from a tool result OK   visible state change OK   failure + fallback OK   clear end condition OK", 0, {"bold": True, "color": NAVY, "size": 15}),
], y=6.5, h=0.6)

# ------------------------------------------- slide 5: mock-LLM seams ---
s = slide("Where AI Lives 1 real ML component, 3 mock-LLM seams",
          "AI is allowed to read, interpret, and advise deterministic code decides and acts")
seams = [
    ("RAG Semantic Retriever REAL ML", TEAL,
     "sentence-transformers all-MiniLM-L6-v2 · cosine similarity over 3 knowledge docs",
     "Triggered only when confidence < 0.70. Purely advisory: it backs the re-plan decision, it never makes it."),
    ("NLPMissionParser mock LLM (seam 1)", ORANGE,
     "Extracts target, capabilities, constraints from free text",
     "Same contract an LLM call would fill. Proposes only capabilities can be added, never removed; pre-flight still gates them."),
    ("FailureAdvisor mock LLM (seam 2)", ORANGE,
     "Classifies tool failures: transient vs permanent",
     "Only selects between pre-approved policies (retry vs abort). The retry budget (3) is still guardrail-enforced."),
    ("Report narrative writer mock LLM (seam 3)", ORANGE,
     "Turns structured mission facts into prose",
     "Can reword, never change: finding, confidence history, and outcome come from mission memory, not the writer."),
]
y = 1.35
for title, c, what, boundary in seams:
    box(s, 0.5, y, 3.4, 1.15, title, c, WHITE, 12)
    tb = s.shapes.add_textbox(Inches(4.15), y + 0.02, Inches(8.6), Inches(1.15))
    tf = tb.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; r1 = p1.add_run(); r1.text = what
    r1.font.size = Pt(13); r1.font.bold = True; r1.font.name = MONO; r1.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = boundary
    r2.font.size = Pt(12); r2.font.color.rgb = GRAY
    y += 1.28
box(s, 0.5, 6.6, 12.3, 0.65,
    "A real model drops into any seam tomorrow the guardrails never trust it either way.",
    NAVY, WHITE, 14)

# ---------------------------------------- slide 6: failure and safety ---
s = slide("Failure Handling & Safety", "Detect > interpret > fallback > update state > continue or stop safely")
bullets(s, [
    ("Two planted failures, two different fallbacks", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("Camera timeout > advisor says transient > bounded retry (3).   Weak confidence > not a failure > more evidence (re-inspect).", 0, {"size": 15}),
    ("Where I deliberately do NOT trust the agent", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("Pre-flight gate: 5 critical checks must pass before any movement (low-battery scenario > abort with zero movement commands)", 0, {"size": 15}),
    ("Battery floor 30% (or a parsed free-text floor) for movement · capture retry cap · no capture at base · no detection without evidence · 30-step budget", 0, {"size": 15}),
    ("Defense in depth: even past the guardrail, MockDrone itself refuses moves <= 20% battery like a real flight controller.", 0, {"size": 15}),
    ("Every tool call is watchdogged", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("A hang > 10s becomes a failure ToolObservation and flows through the normal failure pipeline guardrails still run. (Honest limit: the stuck thread keeps running; v2 adds process isolation.)", 0, {"size": 15}),
    ("End conditions are always defined", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("COMPLETE with report · ABORTED with printed reason + failure report · step budget exhausted > abort. No infinite loops.", 0, {"size": 15}),
], y=1.3, h=4.3)
box(s, 0.6, 5.85, 12.1, 1.0,
    "Governing principle:  the agent proposes what to do  deterministic software decides whether it is allowed.",
    NAVY, WHITE, 17)

# ------------------------------------------------------ slide 7: ownership ---
s = slide("My Ownership: the Decision Core",
          "AgentController._decide() + MissionState small enough to defend line by line")
bullets(s, [
    ("Internals", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("_decide() = pure function of MissionState > AgentDecision (a WHY), which expands into validated AgentActions (a WHAT).", 0, {"size": 15}),
    ("Decisions and actions are different types: the validator refuses anything that is not an executable action no path from 're-inspect' to the drone without becoming a checked command.", 0, {"size": 15}),
    ("What can fail and how it is handled", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("Tool failure > observation, not exception. Guardrail rejection > recorded SAFETY_REJECTION + abort. Retrieval down > caught, mission continues. No convergence > step budget abort.", 0, {"size": 15}),
    ("How I debug it", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("Deterministic scenarios (same input = same trace) · full event history in MissionState · state-in/decision-out makes _decide() hand-testable the unit tests do exactly that.", 0, {"size": 15}),
    ("Alternatives I considered", 0, {"bold": True, "color": NAVY, "size": 18}),
    ("LLM in the loop (cannot defend each decision) · agent framework (hides the state machine I must own) · plan-index script (cannot respond to change). Honest limits: one mission type, one-if failure classifier, tool timeout lacks process isolation v2 roadmap has the fixes.", 0, {"size": 15}),
], y=1.3, h=5.9, size=15)

prs.save("WeevilDrone_Agentic_Demo.pptx")
print("saved WeevilDrone_Agentic_Demo.pptx, slides:", len(prs.slides._sldIdLst))
